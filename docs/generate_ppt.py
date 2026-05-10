
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG    = RGBColor(0x0D, 0x0D, 0x0D)
ORANGE     = RGBColor(0xF9, 0x73, 0x16)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x9C, 0xA3, 0xAF)
CARD_BG    = RGBColor(0x1F, 0x29, 0x37)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
RED        = RGBColor(0xEF, 0x44, 0x44)

BLANK = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txb(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def header_bar(slide, title="PlateUp 🍽️"):
    rect(slide, 0, 0, 13.33, 0.6, RGBColor(0x1F, 0x29, 0x37))
    txb(slide, "PlateUp 🍽️", 0.2, 0.08, 4, 0.5, size=16, bold=True, color=ORANGE)
    txb(slide, title, 4, 0.08, 9, 0.5, size=16, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, title, body, icon="", title_color=ORANGE):
    rect(slide, l, t, w, h, CARD_BG)
    txb(slide, icon + " " + title, l+0.1, t+0.08, w-0.2, 0.4, size=14, bold=True, color=title_color)
    txb(slide, body, l+0.1, t+0.5, w-0.2, h-0.6, size=12, color=WHITE)


# ── SLIDE 1: TITLE ──────────────────────────────────────────
s1 = add_slide()
bg(s1)
rect(s1, 0, 2.5, 13.33, 0.06, ORANGE)
txb(s1, "PlateUp", 0, 1.0, 13.33, 1.5, size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s1, "Your AI Waiter, Always Ready 🍽️", 0, 2.6, 13.33, 0.8, size=24, color=ORANGE, align=PP_ALIGN.CENTER)
txb(s1, "Shubham Bhandari  |  AI Fellowship  |  May 4, 2026", 0, 6.5, 13.33, 0.7, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 2: PROBLEM STATEMENT ──────────────────────────────
s2 = add_slide()
bg(s2)
header_bar(s2, "Problem Statement")
txb(s2, '"Every minute a waiter is busy, a customer is frustrated."', 0.3, 0.7, 12.7, 0.7, size=20, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)
problems = [
    ("🕐", "Long Wait Times",       "Customers wait too long for a waiter"),
    ("❌", "Manual Order Errors",   "Wrong food delivered, guests upset"),
    ("🌙", "No 24/7 Support",       "Orders lost outside business hours"),
    ("📡", "No Real-Time Updates",  "Customers unaware of order progress"),
    ("📋", "Manual Reporting",      "Managers have no data insights"),
    ("🤖", "No Personalization",    "Generic service, no recommendations"),
]
cols = [(0.2, 1.55), (4.6, 1.55), (9.0, 1.55),
        (0.2, 4.0),  (4.6, 4.0),  (9.0, 4.0)]
for i, (icon, title, body) in enumerate(problems):
    x, y = cols[i]
    card(s2, x, y, 4.1, 2.2, title, body, icon)

# ── SLIDE 3: SOLUTION OVERVIEW ──────────────────────────────
s3 = add_slide()
bg(s3)
header_bar(s3, "Solution Overview")
txb(s3, "Meet PlateUp — Agentic AI for Restaurants", 0, 0.65, 13.33, 0.6, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
boxes = [
    (0.3,  1.5, "Customer\nMessage",    CARD_BG),
    (2.5,  1.5, "Chat\nInterface",      CARD_BG),
    (4.7,  1.5, "ChatServlet\n(Java)",  CARD_BG),
]
for l, t, txt, clr in boxes:
    rect(s3, l, t, 1.9, 1.0, clr)
    txb(s3, txt, l, t+0.1, 1.9, 0.8, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s3, "→", l+2.0, t+0.2, 0.4, 0.6, size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

branches = [
    (7.2, 1.3, "Known Command\n→ OrderService → MySQL",          "🗃️"),
    (7.2, 2.6, "Natural Language\n→ Gemini 2.0 Flash → AI Reply","🤖"),
    (7.2, 3.9, "Order Confirmed\n→ n8n → Email Automation",      "📧"),
]
for l, t, txt, icon in branches:
    rect(s3, l, t, 5.8, 1.0, CARD_BG)
    txb(s3, icon + "  " + txt, l+0.15, t+0.1, 5.5, 0.85, size=13, color=WHITE)
txb(s3, "One AI brain. Three powerful flows.", 0, 6.7, 13.33, 0.5, size=16, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ── SLIDE 4: KEY FEATURES ───────────────────────────────────
s4 = add_slide()
bg(s4)
header_bar(s4, "Key Features")
features = [
    ("🤖", "Agentic AI Chat",         "Natural language via Gemini 2.0 Flash"),
    ("📧", "Email Automation",         "Order confirmation email sent via n8n + Gmail"),
    ("🍽️", "Smart Menu Browsing",     "Filter by veg, spicy, or price"),
    ("📦", "Real-Time Order Tracking", "4-stage: Placed→Preparing→Ready→Served"),
    ("👨‍💼","Staff Dashboard",          "Live order queue with status control"),
    ("🌙", "Dark / Light Theme",       "Persistent user-preferred UI"),
    ("💳", "Seamless Checkout",        "Cart auto-clears after payment"),
]
positions = [
    (0.2, 0.8), (4.77, 0.8), (9.33, 0.8),
    (0.2, 3.3), (4.77, 3.3), (9.33, 3.3),
    (4.77, 5.8),
]
for i, (icon, title, body) in enumerate(features):
    x, y = positions[i]
    card(s4, x, y, 4.3, 2.2, title, body, icon)

# ── SLIDE 5: TECH STACK ─────────────────────────────────────
s5 = add_slide()
bg(s5)
header_bar(s5, "Tech Stack")
txb(s5, "Technologies Used", 0, 0.65, 13.33, 0.5, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
stack = [
    ("Backend",    "Java 21 · Jakarta EE 6.0 · Apache Tomcat 10.1"),
    ("AI Engine",  "Google Gemini 2.0 Flash"),
    ("Automation", "n8n 2.18.5  ·  Gmail Node (Email)"),
    ("Database",   "MySQL 8.0 + JDBC Connector/J 8.3.0"),
    ("Frontend",   "JSP · HTML5 · CSS3 · Vanilla JS"),
    ("Build/VCS",  "Maven 3.x · Git · GitHub"),
]
for i, (layer, tech) in enumerate(stack):
    y = 1.3 + i * 0.92
    rect(s5, 0.5, y, 3.0, 0.75, CARD_BG)
    txb(s5, layer, 0.55, y+0.1, 2.9, 0.55, size=14, bold=True, color=ORANGE)
    rect(s5, 3.7, y, 9.0, 0.75, RGBColor(0x11, 0x18, 0x27))
    txb(s5, tech,  3.85, y+0.1, 8.7, 0.55, size=14, color=WHITE)

# ── SLIDE 6: SYSTEM ARCHITECTURE ────────────────────────────
s6 = add_slide()
bg(s6)
header_bar(s6, "System Architecture")
txb(s6, "4-Layer Architecture with Agentic AI Middleware", 0, 0.65, 13.33, 0.5, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
layers = [
    (2.5, 1.2,  8.3, 0.85, "PRESENTATION LAYER  —  JSP · HTML5 · CSS3 · Vanilla JS"),
    (2.5, 2.3,  8.3, 0.85, "BUSINESS LOGIC LAYER  —  Servlets · Service Classes · DAO"),
    (0.5, 3.4,  5.5, 0.85, "AGENTIC AI LAYER  —  Google Gemini 2.0 Flash API"),
    (7.0, 3.4,  5.5, 0.85, "AUTOMATION LAYER  —  n8n Webhooks · Gmail Email"),
    (2.5, 4.5,  8.3, 0.85, "DATA LAYER  —  MySQL 8.0 · JDBC · DAO Pattern"),
]
for l, t, w, h, txt in layers:
    rect(s6, l, t, w, h, CARD_BG)
    txb(s6, txt, l+0.15, t+0.15, w-0.3, 0.6, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for y in [2.1, 3.2, 4.35]:
    txb(s6, "↓", 6.4, y, 0.5, 0.3, size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

txb(s6, "MVC Pattern:  Model (Java POJOs)  →  View (JSP)  →  Controller (Servlets)", 0, 5.6, 13.33, 0.5, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 7: AGENTIC AI IMPLEMENTATION ──────────────────────
s7 = add_slide()
bg(s7)
header_bar(s7, "Agentic AI — How It Works")
txb(s7, "Perceive  →  Reason  →  Act", 0, 0.65, 13.33, 0.5, size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
steps = [
    ("1. PERCEIVE", "ChatServlet receives any message not matching known commands"),
    ("2. REASON",   "GeminiService builds a prompt with full menu context, prices, and rules"),
    ("3. ACT",      "Gemini 2.0 Flash returns a natural, helpful reply with dish names & prices"),
]
for i, (title, body) in enumerate(steps):
    card(s7, 0.3, 1.3 + i*1.5, 5.5, 1.3, title, body)

txb(s7, "BEFORE (Keyword Matching)", 6.3, 1.1, 6.6, 0.4, size=13, bold=True, color=RED)
before_code = 'if (message.contains("veg")) {\n    response = showVegMenu();\n}\n// Fails: "no meat please" → broken'
rect(s7, 6.3, 1.5, 6.6, 1.3, RGBColor(0x1a, 0x00, 0x00))
txb(s7, before_code, 6.45, 1.55, 6.3, 1.2, size=11, color=WHITE)

txb(s7, "AFTER (Agentic AI)", 6.3, 3.1, 6.6, 0.4, size=13, bold=True, color=GREEN)
after_code = 'String reply = geminiService\n  .getAIResponse(message, allMenuItems);\n// Works on ANY natural language ✅'
rect(s7, 6.3, 3.5, 6.6, 1.2, RGBColor(0x00, 0x1a, 0x00))
txb(s7, after_code, 6.45, 3.55, 6.3, 1.1, size=11, color=WHITE)

rect(s7, 0.3, 5.6, 12.7, 1.3, CARD_BG)
txb(s7, '💬 Example:  Customer says: "I want something spicy and cheap"\nPlateUp AI: "Try Chicken Tikka (₹280, Spice 4/5) 🌶️ or Prawn Masala (₹440)! Type: add 1 chicken tikka"', 0.5, 5.7, 12.3, 1.1, size=13, color=WHITE)

# ── SLIDE 8: n8n + EMAIL AUTOMATION ─────────────────────────
s8 = add_slide()
bg(s8)
header_bar(s8, "n8n + Email Automation")
txb(s8, "Automated Email Confirmation on Every Order", 0, 0.65, 13.33, 0.5, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

flow = [
    ("📦", "Order\nConfirmed"),
    ("☕", "Java Fires\nWebhook"),
    ("⚙️", "n8n\nReceives"),
    ("📧", "Gmail Node\nSends Email"),
]
for i, (icon, label) in enumerate(flow):
    x = 0.5 + i * 3.1
    rect(s8, x, 1.2, 2.6, 1.2, CARD_BG)
    txb(s8, icon, x, 1.3, 2.6, 0.5, size=22, align=PP_ALIGN.CENTER)
    txb(s8, label, x, 1.75, 2.6, 0.6, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        txb(s8, "→", x+2.65, 1.55, 0.4, 0.5, size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

json_txt = '{\n  "orderId": 42,\n  "customerName": "Shubham",\n  "customerEmail": "shubham@test.com",\n  "items": "2x Butter Chicken",\n  "timestamp": "2026-05-04T..."\n}'
rect(s8, 0.3, 2.7, 6.0, 3.5, RGBColor(0x11, 0x18, 0x27))
txb(s8, "📤 Payload sent to n8n:", 0.4, 2.75, 5.8, 0.4, size=13, bold=True, color=ORANGE)
txb(s8, json_txt, 0.5, 3.15, 5.6, 2.9, size=12, color=WHITE)

props = [
    ("✅ Automatic",    "No human click needed — fires on order"),
    ("✅ Real Email",   "Customer receives Gmail confirmation"),
    ("✅ Extendable",   "Add WhatsApp / Slack in minutes"),
    ("✅ Observable",   "n8n shows full execution history"),
]
rect(s8, 6.7, 2.7, 6.3, 3.5, CARD_BG)
txb(s8, "Why n8n Automation?", 6.8, 2.8, 6.0, 0.4, size=14, bold=True, color=ORANGE)
for i, (title, body) in enumerate(props):
    txb(s8, title, 6.9, 3.3 + i*0.7, 2.0, 0.5, size=13, bold=True, color=GREEN)
    txb(s8, body,  9.0, 3.3 + i*0.7, 3.8, 0.5, size=12, color=WHITE)

# ── SLIDE 9: CHALLENGES & SOLUTIONS ─────────────────────────
s9 = add_slide()
bg(s9)
header_bar(s9, "Challenges & Solutions")
txb(s9, "Real Problems. Real Fixes.", 0, 0.65, 13.33, 0.5, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
challenges = [
    ("Gemini 1.5-flash deprecated",       "Migrated to gemini-2.0-flash — worked instantly"),
    ("n8n IF node type mismatch",          "Removed IF node — simplified to 3-node pipeline"),
    ("Served orders showing in cart",      "Filter logic updated to skip status = 'served'"),
    ("n8n install failed on Windows",      "Used npm install -g n8n — solved immediately"),
    ("isVeg() method not found",           "Corrected to isVegetarian() per MenuItem model"),
    ("API key risk in Git commits",        "Added .gitignore — key rotation planned post-demo"),
]
for i, (prob, sol) in enumerate(challenges):
    row = i // 2
    col = i % 2
    x = 0.3 + col * 6.55
    y = 1.3 + row * 1.9
    rect(s9, x, y, 6.2, 1.7, CARD_BG)
    txb(s9, "❌ " + prob, x+0.15, y+0.1,  5.9, 0.6, size=12, bold=True, color=RED)
    txb(s9, "✅ " + sol,  x+0.15, y+0.7,  5.9, 0.8, size=12, color=GREEN)

# ── SLIDE 10: RESULTS ───────────────────────────────────────
s10 = add_slide()
bg(s10)
header_bar(s10, "Results & What Works")
txb(s10, "Everything is Live and Working ✅", 0, 0.65, 13.33, 0.5, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
stats = [
    ("< 3 sec",   "Gemini AI\nResponse Time"),
    ("< 500 ms",  "n8n Webhook\nResponse"),
    ("< 1 sec",   "Page\nLoad Time"),
    ("< 100 ms",  "Database\nQuery Time"),
]
for i, (num, label) in enumerate(stats):
    x = 0.5 + i * 3.1
    rect(s10, x, 1.2, 2.7, 1.5, CARD_BG)
    txb(s10, num,   x, 1.3,  2.7, 0.7, size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txb(s10, label, x, 1.95, 2.7, 0.65, size=12, color=GRAY, align=PP_ALIGN.CENTER)

checks = [
    "✅ User Login & Registration (Customer + Staff)",
    "✅ Gemini 2.0 Flash AI Chat",
    "✅ Smart Menu Browsing with Filters",
    "✅ Real-Time 4-Stage Order Tracking",
    "✅ Automated Email Confirmation via n8n + Gmail",
    "✅ Staff Dashboard with Live Order Queue",
    "✅ Cart Auto-Clears After Payment",
    "✅ Dark / Light Theme + Responsive Design",
]
half = len(checks) // 2
for i, c in enumerate(checks[:half]):
    txb(s10, c, 0.4, 2.9 + i*0.7, 6.3, 0.6, size=13, color=WHITE)
for i, c in enumerate(checks[half:]):
    txb(s10, c, 6.9, 2.9 + i*0.7, 6.3, 0.6, size=13, color=WHITE)

# ── SLIDE 11: FUTURE SCOPE ──────────────────────────────────
s11 = add_slide()
bg(s11)
header_bar(s11, "Future Scope")
txb(s11, "What's Next for PlateUp 🚀", 0, 0.65, 13.33, 0.5, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
future = [
    ("📱", "WhatsApp Notifications",
     "Automatically send WhatsApp messages to the restaurant owner and customer when an order is placed — using Twilio + n8n, no extra code needed."),
    ("💳", "Payment Interface",
     "Integrate a real payment gateway like Razorpay or UPI so customers can pay online directly from the app instead of paying manually at the table."),
]
for i, (icon, title, body) in enumerate(future):
    x = 0.5 + i * 6.4
    rect(s11, x, 1.4, 6.0, 4.5, CARD_BG)
    txb(s11, icon, x, 1.55, 6.0, 0.9, size=48, align=PP_ALIGN.CENTER)
    txb(s11, title, x, 2.5, 6.0, 0.6, size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txb(s11, body, x+0.2, 3.2, 5.6, 2.5, size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ── SLIDE 12: THANK YOU ─────────────────────────────────────
s12 = add_slide()
bg(s12)
rect(s12, 0, 3.1, 13.33, 0.06, ORANGE)
txb(s12, "Thank You 🙏", 0, 0.8, 13.33, 1.2, size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s12, "PlateUp — Where AI Meets Hospitality", 0, 2.1, 13.33, 0.6, size=22, color=ORANGE, align=PP_ALIGN.CENTER)
takeaways = [
    ("🤖", "Agentic AI",      "Gemini 2.0 Flash\nPerceive → Reason → Act"),
    ("📧", "Email Automation","n8n + Gmail\nAuto order confirmation"),
    ("🏗️", "Clean Code",     "MVC + DAO\nRole-Based Auth"),
    ("👨‍💻","Student Project", "AI Fellowship\nMay 4, 2026"),
]
for i, (icon, title, body) in enumerate(takeaways):
    x = 0.4 + i * 3.15
    rect(s12, x, 3.4, 2.9, 2.0, CARD_BG)
    txb(s12, icon + " " + title, x+0.1, 3.5, 2.7, 0.5, size=13, bold=True, color=ORANGE)
    txb(s12, body, x+0.1, 4.0, 2.7, 1.2, size=12, color=WHITE)

txb(s12, "github.com/shubhambhandari24/Restaurant_chatbot  |  Shubham Bhandari  |  AI Fellowship  |  May 4, 2026",
    0, 6.6, 13.33, 0.6, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ── SAVE ─────────────────────────────────────────────────────
out = r"c:\Users\shubh\eclipse\jee-2026-03\Restaurant_Chatbot\docs\PlateUp_Presentation.pptx"
prs.save(out)
print("Saved:", out)
