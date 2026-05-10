from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLACK  = RGBColor(0x0D,0x0D,0x0D)
DARK   = RGBColor(0x1A,0x1A,0x2E)
CARD   = RGBColor(0x16,0x21,0x3E)
GOLD   = RGBColor(0xF5,0x9E,0x0B)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GREY   = RGBColor(0x9C,0xA3,0xAF)
GREEN  = RGBColor(0x10,0xB9,0x81)
BLUE   = RGBColor(0x3B,0x82,0xF6)
PURPLE = RGBColor(0xA8,0x55,0xF7)
RED    = RGBColor(0xEF,0x44,0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def rect(slide, l, t, w, h, fill=DARK):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return txb

def bg(slide): rect(slide, 0, 0, 13.33, 7.5, DARK)

def header(slide, title, sub=None):
    txt(slide, title, 0.5, 0.15, 12.33, 0.6, size=28, bold=True, color=GOLD)
    rect(slide, 0.5, 0.78, 12.33, 0.06, GOLD)
    if sub:
        txt(slide, sub, 0.5, 0.88, 12.33, 0.38, size=12, color=GREY, italic=True)

def card(slide, l, t, w, h, accent, title, lines):
    rect(slide, l, t, w, h, CARD)
    rect(slide, l, t, w, 0.07, accent)
    txt(slide, title, l+0.15, t+0.1, w-0.3, 0.42, size=13, bold=True, color=accent)
    y = t + 0.58
    for line in lines:
        txt(slide, f"• {line}", l+0.15, y, w-0.3, 0.38, size=10, color=WHITE)
        y += 0.38

# ── SLIDE 1 — TITLE ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.5, 7.5, GOLD)
rect(s, 1.1, 1.5, 2.3, 2.3, CARD)
txt(s, "🍽️", 1.4, 1.65, 1.7, 1.7, size=52, align=PP_ALIGN.CENTER)
txt(s, "PlateUp", 3.8, 1.4, 9, 1.0, size=58, bold=True, color=GOLD)
txt(s, "Restaurant Chatbot", 3.8, 2.5, 9, 0.65, size=26, color=WHITE)
rect(s, 3.8, 3.25, 8.8, 0.07, GOLD)
txt(s, "Powered by Agentic AI  ·  Google Gemini 2.0  ·  n8n Automation", 3.8, 3.38, 9, 0.45, size=13, color=GREY, italic=True)
for i,(tag,col) in enumerate([("Java EE",BLUE),("Gemini AI",PURPLE),("n8n Cloud",GREEN),("MySQL",GOLD)]):
    rect(s, 3.8+i*2.1, 4.05, 1.85, 0.44, col)
    txt(s, tag, 3.8+i*2.1, 4.05, 1.85, 0.44, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "AI Fellowship Internship Report  |  Shubham Bhandari  |  May 2026", 3.8, 6.7, 9, 0.45, size=12, color=GREY, italic=True)

# ── SLIDE 2 — PROBLEM STATEMENT ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Problem Statement", "Why does the restaurant industry need an AI-powered chatbot?")
problems = [
    ("⏳ Long Wait Times",    "Customers wait for waitstaff to take orders — causing delays and frustration."),
    ("📋 Order Errors",       "Manual note-taking leads to wrong food delivered and dissatisfied customers."),
    ("🕐 No 24/7 Support",    "Restaurants can't afford round-the-clock staff for every customer query."),
    ("📊 No Real-time Data",  "Managers lack live order analytics, peak-hour insights, and demand forecasting."),
]
for i,(title,desc) in enumerate(problems):
    c,r = i%2, i//2
    l,t = 0.5+c*6.45, 1.38+r*2.6
    rect(s, l, t, 6.1, 2.3, CARD)
    rect(s, l, t, 6.1, 0.07, GOLD)
    txt(s, title, l+0.2, t+0.12, 5.7, 0.45, size=14, bold=True, color=GOLD)
    txt(s, desc,  l+0.2, t+0.65, 5.7, 1.45, size=11, color=WHITE)
txt(s, "✅  PlateUp solves ALL of these with Agentic AI + Automation", 0.5, 6.8, 12.33, 0.4, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ── SLIDE 3 — SOLUTION OVERVIEW ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Solution Overview", "4-layer architecture with Agentic AI at the core")
layers = [
    (BLUE,   "🌐  Presentation",   "JSP · HTML5 · CSS3 · JavaScript",           0.5),
    (PURPLE, "⚙️   Business Logic", "Java Servlets · Services · DAO Pattern",     3.45),
    (GREEN,  "🤖  Agentic AI",     "Gemini 2.0 Flash · n8n Cloud Webhooks",      6.4),
    (GOLD,   "🗄️   Data Layer",    "MySQL 8.0 · JDBC · Connection Pool",         9.35),
]
for col,title,desc,l in layers:
    rect(s, l, 1.38, 2.7, 4.9, CARD)
    rect(s, l, 1.38, 2.7, 0.08, col)
    txt(s, title, l+0.15, 1.52, 2.4, 0.5,  size=12, bold=True, color=col)
    txt(s, desc,  l+0.15, 2.1,  2.4, 3.8,  size=10, color=WHITE)
for ax in [3.2,6.15,9.1]:
    txt(s, "→", ax, 3.5, 0.3, 0.45, size=22, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "Apache Tomcat 10.1  ·  Maven 3.x  ·  Eclipse IDE 2026-03  ·  Git + GitHub", 0.5, 6.6, 12.33, 0.45, size=11, color=GREY, italic=True, align=PP_ALIGN.CENTER)

# ── SLIDE 4 — AGENTIC AI ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Agentic AI — Google Gemini 2.0 Flash", "Perceive → Reason → Act: beyond simple keyword matching")
rect(s, 0.5, 1.3, 12.33, 1.15, CARD)
rect(s, 0.5, 1.3, 12.33, 0.08, PURPLE)
txt(s, '💬 Customer: "What do you recommend for a spicy vegetarian dinner under ₹200?"', 0.7, 1.4, 11.9, 0.42, size=12, bold=True, color=WHITE)
txt(s, '🤖 PlateUp AI: "For spicy veggie under ₹200, try Dal Makhani (₹180, Spice 4/5) 🌶️ or Mushroom Masala (₹160)! Both are crowd favourites. Type: add 1 dal makhani"', 0.7, 1.84, 11.9, 0.52, size=11, color=GREEN)
pillars = [
    (PURPLE,"🧠  Perceive", "Receives ANY natural language input. Understands intent — not hardcoded keywords."),
    (BLUE,  "⚙️  Reason",  "Gemini 2.0 analyses full menu: price, category, spice level, veg/non-veg flag."),
    (GREEN, "🎯  Act",     "Generates helpful, contextual reply. Guides customer to place order using commands."),
]
for i,(col,title,desc) in enumerate(pillars):
    l = 0.5+i*4.15
    rect(s, l, 2.62, 3.9, 3.25, CARD)
    rect(s, l, 2.62, 3.9, 0.08, col)
    txt(s, title, l+0.2, 2.75, 3.5, 0.45, size=14, bold=True, color=col)
    txt(s, desc,  l+0.2, 3.3,  3.5, 2.3,  size=11, color=WHITE)
txt(s, "Model: gemini-2.0-flash  ·  Free: 1,500 req/day  ·  Max tokens: 150  ·  Temperature: 0.7  ·  Zero cost", 0.5, 6.1, 12.33, 0.38, size=11, color=GREY, italic=True, align=PP_ALIGN.CENTER)

# ── SLIDE 5 — n8n AUTOMATION ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "n8n Agentic Automation Layer", "Live cloud workflow — fires automatically on every order")
flow = [
    (BLUE,   "☕ Java App",       "ChatServlet fires POST webhook after item added to order"),
    (PURPLE, "🔗 n8n Cloud",      "shubham2408.app.n8n.cloud receives payload in real-time"),
    (GOLD,   "📧 Gmail Node",     "Auto-sends personalized order confirmation to customer email"),
    (GREEN,  "✅ Execution Log",  "Every run logged with timestamp, status, order data"),
]
for i,(col,title,desc) in enumerate(flow):
    l = 0.4+i*3.1
    rect(s, l, 1.45, 2.85, 2.1, CARD)
    rect(s, l, 1.45, 2.85, 0.08, col)
    txt(s, title, l+0.15, 1.58, 2.55, 0.45, size=12, bold=True, color=col)
    txt(s, desc,  l+0.15, 2.1,  2.55, 1.3,  size=10, color=WHITE)
    if i < 3:
        txt(s, "→", l+2.87, 2.28, 0.35, 0.4, size=22, color=GOLD, align=PP_ALIGN.CENTER)
rect(s, 0.5, 3.75, 12.33, 2.55, CARD)
rect(s, 0.5, 3.75, 12.33, 0.07, GREEN)
txt(s, "✅ What the Email Contains (Live Demo)", 0.7, 3.86, 11.9, 0.42, size=13, bold=True, color=GREEN)
email_lines = [
    "📬 To: {{ customerEmail }} — dynamically filled from order data",
    "📌 Subject: 🍽️ Order Confirmed - PlateUp Restaurant",
    "📝 Body: Customer name, Order ID, Items ordered, Total amount — all auto-filled",
    "🚀 Delivery: <500ms after order is placed — zero human intervention",
]
for i,line in enumerate(email_lines):
    txt(s, line, 0.7, 4.38+i*0.48, 11.9, 0.42, size=11, color=WHITE)

# ── SLIDE 6 — KEY FEATURES ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Key Features", "End-to-end restaurant intelligence with AI at the core")
features = [
    (GOLD,   "🤖 AI Chat Ordering",    ["Gemini 2.0 natural language","30+ context-aware intents","Fuzzy typo correction","Smart dish cards after AI reply"]),
    (BLUE,   "📋 Order Management",    ["Chat + button ordering","4-stage live status tracker","Auto cart clear after payment","Live cart badge counter"]),
    (GREEN,  "📧 n8n Email Automation",["Auto Gmail on every order","Dynamic customer data fill","Cloud-hosted — always on","Execution log for proof"]),
    (PURPLE, "👤 User Profile",        ["Edit name & phone","Secure password change","Order history view","Food preferences (localStorage)"]),
    (GOLD,   "🍽️ Rich Menu UI",        ["Visual cards with images","Calorie badges per dish","FSSAI veg/non-veg dot","Chinese & Beverages sections"]),
    (BLUE,   "🎨 Premium UI/UX",       ["Dark mode default","Outfit font + gold palette","Responsive + micro-animations","Staff dashboard"]),
]
for i,(col,title,bullets) in enumerate(features):
    r,c = i//3, i%3
    card(s, 0.5+c*4.28, 1.32+r*2.78, 3.98, 2.55, col, title, bullets)

# ── SLIDE 7 — TECH STACK ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Technology Stack", "Industry-standard, production-grade tools — 100% free")
stack = [
    (BLUE,   "☕ Java EE 21",       ["Jakarta Servlet 6.0 API","Apache Tomcat 10.1.54","Maven build tool","Eclipse IDE 2026-03"]),
    (PURPLE, "🤖 Google Gemini AI", ["gemini-2.0-flash model","1,500 free requests/day","REST via HttpURLConnection","Max 150 tokens/response"]),
    (GREEN,  "🔗 n8n Automation",   ["Cloud: shubham2408.app.n8n.cloud","Webhook → Gmail → Log","Visual drag-and-drop builder","Zero extra Java code"]),
    (GOLD,   "🗄️ MySQL 8.0",        ["MySQL Connector/J 8.3.0","JDBC connection pooling","DAO abstraction pattern","5 normalized tables"]),
    (BLUE,   "🎨 Frontend",         ["JSP + HTML5 + CSS3","Vanilla JavaScript","Outfit Google Font","No heavy frameworks"]),
    (PURPLE, "📦 Libraries",        ["org.json 20240303","Jakarta Servlet API","Jakarta JSP API 3.1","Git + GitHub VCS"]),
]
for i,(col,title,bullets) in enumerate(stack):
    r,c = i//3, i%3
    card(s, 0.5+c*4.28, 1.32+r*2.78, 3.98, 2.55, col, title, bullets)

# ── SLIDE 8 — IMPLEMENTATION FLOW ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Key Implementation — Agentic AI Flow", "How a customer message becomes an AI-powered action")
steps = [
    (BLUE,   "1","Customer Types",   "Natural language message in chat UI"),
    (PURPLE, "2","Intent Engine",    "chat.js detectIntent() — 30+ regex patterns"),
    (GOLD,   "3","Fuzzy Match",      "Levenshtein correction for misspelled item names"),
    (GREEN,  "4","GeminiService",    "Unknown intent → builds prompt with menu context"),
    (BLUE,   "5","Gemini 2.0 API",  "POST to Google API → AI generates smart reply"),
    (PURPLE, "6","Smart Cards",      "findAndShowMentionedDishes() renders visual cards"),
    (GOLD,   "7","n8n Webhook",      "Item added → cloud notification → Gmail email"),
    (GREEN,  "8","UI Response",      "Reply shown + cart badge + order status live"),
]
for i,(col,num,title,desc) in enumerate(steps):
    r,c = i//4, i%4
    l,t = 0.4+c*3.15, 1.32+r*2.75
    rect(s, l, t, 2.95, 2.4, CARD)
    rect(s, l, t, 0.5, 2.4, col)
    txt(s, num,   l,     t+0.88, 0.5,  0.45, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, l+0.6, t+0.1,  2.25, 0.45, size=12, bold=True, color=col)
    txt(s, desc,  l+0.6, t+0.62, 2.25, 1.6,  size=10, color=WHITE)

# ── SLIDE 8b — FUZZY MATCH + 30 INTENTS ─────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Fuzzy Matching + 30 Smart Intents", "Client-side intelligence — zero server calls for instant responses")
rect(s, 0.5, 1.3, 6.0, 4.7, CARD)
rect(s, 0.5, 1.3, 6.0, 0.08, PURPLE)
txt(s, "🔍 Levenshtein Fuzzy Matching", 0.7, 1.42, 5.6, 0.45, size=13, bold=True, color=PURPLE)
fuzzy_lines = [
    "Customer types: add pannr tika",
    "→ Levenshtein distance scored against all 30 items",
    "→ Best match: Paneer Tikka (score ≤ 2.5 threshold)",
    "→ Bot: 'Did you mean Paneer Tikka? Adding it! 😊'",
    "→ Item added — no re-typing required",
    "",
    "Handles: paner / tikaa / butr chikn / dal makni",
    "Threshold tuned to avoid false positives",
]
for i,line in enumerate(fuzzy_lines):
    txt(s, line, 0.7, 1.95+i*0.44, 5.6, 0.4, size=10, color=WHITE if line else WHITE)
rect(s, 6.7, 1.3, 6.1, 4.7, CARD)
rect(s, 6.7, 1.3, 6.1, 0.08, GOLD)
txt(s, "⚡ 30+ Context-Aware Intents", 6.9, 1.42, 5.7, 0.45, size=13, bold=True, color=GOLD)
intents = [
    ("☀️ weather_hot",   "Cold drinks & light bites"),
    ("🌧️ weather_cold",  "Warm comfort food picks"),
    ("🌶️ spicy_food",    "Highest spice-level dishes"),
    ("🎉 party",         "Curated party platter"),
    ("⚡ quick",         "Fastest-to-serve items"),
    ("🥗 healthy",       "Dishes ≤ 300 kcal"),
    ("💰 budget",        "Items under ₹200"),
    ("⏰ time_based",    "Uses real clock hour"),
    ("⚠️ allergy",       "Allergy-safe menu subset"),
    ("😮 tired",         "Comfort food selection"),
]
for i,(tag,desc) in enumerate(intents):
    txt(s, f"{tag} — {desc}", 6.9, 1.95+i*0.44, 5.7, 0.4, size=10, color=WHITE)
txt(s, "All handled client-side — instant 0ms response, no server needed", 0.5, 6.2, 12.33, 0.4, size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ── SLIDE 9 — CHALLENGES & SOLUTIONS ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Challenges & Solutions", "Real problems encountered and how they were resolved")
challenges = [
    (RED,   GREEN, "Gemini model deprecated",        "gemini-1.5-flash sunset by Google",          "Migrated to gemini-2.0-flash — newer, faster, same free tier"),
    (RED,   GREEN, "n8n install failed on Windows",  "npx n8n crashed — C++ build tools missing",  "Switched to n8n Cloud — more reliable, no install needed"),
    (RED,   GREEN, "Served orders in cart",           "Paid orders re-appeared in order view",      "Added status filter — skip 'served' orders in query"),
    (RED,   GREEN, "n8n IF node type mismatch",       "Number vs string comparison always failed",  "Removed IF node — simplified to 3-node pipeline"),
    (RED,   GREEN, "isVeg() method not found",        "Incorrect method name caused compile error",  "Corrected to isVegetarian() per MenuItem model"),
    (RED,   GREEN, "API key exposed in code",         "Gemini key hardcoded — security risk",       "Added to .gitignore — key rotation planned post-demo"),
]
for i,(rc,gc,challenge,prob,sol) in enumerate(challenges):
    r,c = i//2, i%2
    l,t = 0.5+c*6.45, 1.32+r*1.8
    rect(s, l, t, 6.1, 1.65, CARD)
    txt(s, f"⚠️  {challenge}", l+0.15, t+0.08, 5.8, 0.38, size=12, bold=True, color=GOLD)
    txt(s, f"Problem: {prob}", l+0.15, t+0.52, 5.8, 0.38, size=10, color=RED)
    txt(s, f"✅ Fix: {sol}",   l+0.15, t+0.95, 5.8, 0.55, size=10, color=GREEN)

# ── SLIDE 10 — RESULTS & IMPACT ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Results & Impact", "Measurable outcomes from the live PlateUp system")
metrics = [
    (BLUE,   "< 100ms",  "Database Query Time"),
    (PURPLE, "< 500ms",  "n8n Webhook Response"),
    (GREEN,  "2-3 sec",  "Gemini AI Response"),
    (GOLD,   "< 1 sec",  "Page Load Time"),
]
for i,(col,val,label) in enumerate(metrics):
    l = 0.5+i*3.1
    rect(s, l, 1.35, 2.85, 2.0, CARD)
    rect(s, l, 1.35, 2.85, 0.08, col)
    txt(s, val,   l+0.15, 1.55, 2.55, 0.72, size=28, bold=True, color=col,   align=PP_ALIGN.CENTER)
    txt(s, label, l+0.15, 2.35, 2.55, 0.72, size=11,             color=WHITE, align=PP_ALIGN.CENTER)
working = [
    "✅ Fuzzy typo correction — pannr tika → Paneer Tikka",
    "✅ 30+ smart intents (weather, mood, time, allergy)",
    "✅ Visual menu cards with images + calorie badges",
    "✅ Chinese & Beverages menu sections",
    "✅ AI dish cards auto-shown after Gemini reply",
]
working2 = [
    "✅ User Profile — edit details + change password",
    "✅ Gmail auto email via n8n Cloud",
    "✅ Live cart badge counter on nav",
    "✅ Dark mode as default (localStorage persist)",
    "✅ Staff dashboard + 4-stage order tracking",
]
rect(s, 0.5,  3.55, 6.1,  3.55, CARD)
rect(s, 6.75, 3.55, 6.08, 3.55, CARD)
rect(s, 0.5,  3.55, 6.1,  0.07, GREEN)
rect(s, 6.75, 3.55, 6.08, 0.07, GREEN)
txt(s, "What's Working", 0.7, 3.65, 5.7, 0.42, size=13, bold=True, color=GREEN)
txt(s, "What's Working", 6.95,3.65, 5.7, 0.42, size=13, bold=True, color=GREEN)
for i,line in enumerate(working):
    txt(s, line, 0.7, 4.18+i*0.52, 5.7, 0.45, size=10, color=WHITE)
for i,line in enumerate(working2):
    txt(s, line, 6.95,4.18+i*0.52, 5.7, 0.45, size=10, color=WHITE)

# ── SLIDE 11 — FUTURE SCOPE ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Future Scope", "The roadmap from demo project to production SaaS platform")
phases = [
    (GREEN,  "🟢 Short Term (1-3 months)",  ["WhatsApp alerts via n8n + Twilio","Google Sheets order analytics","Voice ordering via Web Speech API","Environment variable for API key security"]),
    (BLUE,   "🔵 Medium Term (3-6 months)", ["Razorpay / UPI payment gateway","AI recommendation engine","Mobile PWA version","Multi-language support"]),
    (PURPLE, "🟣 Long Term (6+ months)",    ["Multi-restaurant SaaS platform","Kitchen Display System (WebSockets)","AI demand forecasting & inventory","Customer loyalty points program"]),
]
for i,(col,title,items) in enumerate(phases):
    card(s, 0.5+i*4.28, 1.32, 3.98, 5.55, col, title, items)
txt(s, "💡  PlateUp is designed to scale — every layer is modular and extensible", 0.5, 7.05, 12.33, 0.38, size=12, bold=True, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

# ── SLIDE 12 — THANK YOU ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.5, 7.5, GOLD)
rect(s, 0, 6.8, 13.33, 0.7, GOLD)
txt(s, "🍽️", 1.2, 0.8, 2.5, 2.5, size=72, align=PP_ALIGN.CENTER)
txt(s, "Thank You!", 4.0, 0.8, 9, 1.1, size=52, bold=True, color=GOLD)
txt(s, "PlateUp — AI-Powered Restaurant Chatbot", 4.0, 1.98, 9, 0.58, size=20, color=WHITE)
rect(s, 4.0, 2.68, 8.8, 0.07, GOLD)
details = [
    ("👤 Student",    "Shubham Bhandari"),
    ("📅 Date",       "May 2026 — AI Fellowship Internship Report"),
    ("💻 GitHub",     "github.com/shubhambhandari24/Restaurant_chatbot"),
    ("🔗 n8n Cloud",  "shubham2408.app.n8n.cloud"),
    ("🤖 AI Engine",  "Google Gemini 2.0 Flash (Free Tier)"),
    ("✨ New Features","Fuzzy Match · 30 Intents · Profile · Email Auto"),
]
for i,(label,val) in enumerate(details):
    txt(s, label, 4.0, 2.85+i*0.62, 2.5,  0.5, size=12, bold=True, color=GOLD)
    txt(s, val,   6.6, 2.85+i*0.62, 6.5,  0.5, size=12, color=WHITE)
txt(s, "Built with ❤️ using Java EE · Google Gemini AI · n8n Automation · MySQL", 0.6, 6.85, 12.1, 0.45, size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────────
out = r"c:\Users\shubh\eclipse\jee-2026-03\Restaurant_Chatbot\docs\PlateUp_Presentation.pptx"
prs.save(out)
print("Saved: " + out)
print("Total slides: " + str(len(prs.slides)))
